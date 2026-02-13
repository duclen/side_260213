"""
다운로드된 파일(XLSX, PDF, HWP, DOCX, PPTX, JPG 등)을
텍스트 또는 CSV로 변환하는 파서 모듈.
변환이 불가능한 파일은 원본 그대로 보존한다.
"""
import csv
from pathlib import Path

import pandas as pd


def parse_file(file_path: Path) -> pd.DataFrame | None:
    """파일 형식에 따라 적절한 파서를 호출하고, DataFrame을 반환한다.
    파싱 불가 시 None을 반환한다."""
    suffix = file_path.suffix.lower()
    try:
        if suffix in (".xlsx", ".xls"):
            return parse_excel(file_path)
        elif suffix == ".pdf":
            return parse_pdf(file_path)
        elif suffix in (".doc", ".docx"):
            return parse_docx(file_path)
        elif suffix in (".ppt", ".pptx"):
            return parse_pptx(file_path)
        elif suffix in (".hwp", ".hwpx"):
            return parse_hwp(file_path)
        elif suffix in (".jpg", ".jpeg", ".png", ".gif"):
            return parse_image(file_path)
        elif suffix == ".txt":
            return parse_txt(file_path)
    except Exception as e:
        print(f"    [파싱 실패] {file_path.name}: {e}")
    return None


def parse_excel(path: Path) -> pd.DataFrame | None:
    """엑셀 파일의 모든 시트를 합쳐 DataFrame으로 반환."""
    xls = pd.ExcelFile(path, engine="openpyxl" if path.suffix == ".xlsx" else None)
    frames = []
    for sheet in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet, header=None)
        df.insert(0, "_sheet", sheet)
        frames.append(df)
    if frames:
        return pd.concat(frames, ignore_index=True)
    return None


def parse_pdf(path: Path) -> pd.DataFrame | None:
    """PDF에서 테이블을 추출한다."""
    import pdfplumber
    rows = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    rows.append(row)
            # 테이블이 없으면 텍스트 추출
            if not tables:
                text = page.extract_text()
                if text:
                    for line in text.split("\n"):
                        rows.append([line.strip()])
    if rows:
        max_cols = max(len(r) for r in rows)
        rows = [r + [""] * (max_cols - len(r)) for r in rows]
        return pd.DataFrame(rows)
    return None


def parse_docx(path: Path) -> pd.DataFrame | None:
    """DOCX에서 테이블과 텍스트를 추출한다."""
    from docx import Document
    doc = Document(str(path))
    rows = []
    # 테이블 추출
    for table in doc.tables:
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
    # 테이블이 없으면 단락 텍스트
    if not rows:
        for para in doc.paragraphs:
            if para.text.strip():
                rows.append([para.text.strip()])
    if rows:
        max_cols = max(len(r) for r in rows)
        rows = [r + [""] * (max_cols - len(r)) for r in rows]
        return pd.DataFrame(rows)
    return None


def parse_pptx(path: Path) -> pd.DataFrame | None:
    """PPTX에서 테이블과 텍스트를 추출한다."""
    from pptx import Presentation
    prs = Presentation(str(path))
    rows = []
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append([f"slide_{slide_num}"] + cells)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        rows.append([f"slide_{slide_num}", text])
    if rows:
        max_cols = max(len(r) for r in rows)
        rows = [r + [""] * (max_cols - len(r)) for r in rows]
        return pd.DataFrame(rows)
    return None


def parse_hwp(path: Path) -> pd.DataFrame | None:
    """HWP 파일에서 텍스트를 추출한다. (olefile 기반)"""
    try:
        import olefile
        if not olefile.isOleFile(str(path)):
            return None
        ole = olefile.OleFileIO(str(path))
        if ole.exists("PrvText"):
            data = ole.openstream("PrvText").read()
            text = data.decode("utf-16-le", errors="ignore")
            rows = [[line.strip()] for line in text.split("\n") if line.strip()]
            if rows:
                return pd.DataFrame(rows, columns=["text"])
        ole.close()
    except ImportError:
        print("    [경고] olefile 미설치 — HWP 파싱 건너뜀")
    except Exception:
        pass
    return None


def parse_image(path: Path) -> pd.DataFrame | None:
    """이미지에서 OCR로 텍스트를 추출한다."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(path)
        text = pytesseract.image_to_string(img, lang="kor+eng")
        rows = [[line.strip()] for line in text.split("\n") if line.strip()]
        if rows:
            return pd.DataFrame(rows, columns=["ocr_text"])
    except ImportError:
        print("    [경고] pytesseract/Pillow 미설치 — OCR 건너뜀")
    except Exception:
        pass
    return None


def parse_txt(path: Path) -> pd.DataFrame | None:
    """텍스트 파일을 DataFrame으로 변환."""
    for enc in ("utf-8", "cp949", "euc-kr"):
        try:
            text = path.read_text(encoding=enc)
            rows = [[line.strip()] for line in text.split("\n") if line.strip()]
            if rows:
                return pd.DataFrame(rows, columns=["text"])
        except (UnicodeDecodeError, UnicodeError):
            continue
    return None


def save_as_csv(df: pd.DataFrame, output_path: Path):
    """DataFrame을 CSV로 저장."""
    df.to_csv(output_path, index=False, encoding="utf-8-sig")